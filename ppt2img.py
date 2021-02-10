from pptx import Presentation
import os


file_path = r'./doc/窄光谱已有报道结构及性能总结.pptx'

prs = Presentation(file_path)

for i, slide in enumerate(prs.slides):
    for obj in slide.shapes:
        try:
            img_data = obj.img.blob
            img_type = obj.img.content_type
            type_key = img_type.find('/') + 1
            img_type = img_type[type_key: ]
            path = 'img/chem_img/'
            if not os.path.exists(path):
                os.mkdir(path)
            image_file = path + str(i) + '.' + img_type
            file = open(image_file, 'wb')
            file.write(image_file)
            file.close()
        except:
            continue
